"""
PPTX Presentation Translator
Translates all text fields in PowerPoint presentations while preserving
images, layout, animations, and all non-text elements.
Based on the Legacy translation engine (core/api_clients legacy_batch_translate).
"""

import os
import re
import time
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from core.api_clients import legacy_batch_translate
from core.config import get_protected_phrases
from core.constants import FINALIZATION_PROGRESS_STEPS
from core.progress import TranslationProgress
from core.run_fusion import fuse_merged_paragraph_runs, merge_adjacent_runs_inplace
from core.text_utils import normalize_run_tags


def _collect_shapes(shapes):
    """Recursively collect all leaf shapes, descending into groups."""
    result = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            result.extend(_collect_shapes(shape.shapes))
        else:
            result.append(shape)
    return result


def _safe_color_key(font):
    """Return a comparable color value, safe against _NoneColor and theme colors."""
    try:
        c = font.color
        if c is None or c.type is None:
            return None
        return str(c.rgb) if c.rgb else None
    except AttributeError:
        return None


def _runs_have_same_font(run_a, run_b):
    """Check if two pptx runs share identical font properties."""
    fa, fb = run_a.font, run_b.font
    return (
        fa.bold == fb.bold
        and fa.italic == fb.italic
        and fa.size == fb.size
        and fa.underline == fb.underline
        and _safe_color_key(fa) == _safe_color_key(fb)
        and (fa.name or "") == (fb.name or "")
    )


def _can_merge_adjacent_pptx_runs(prev, run):
    return bool(run.text) and _runs_have_same_font(prev, run)


def _fuse_paragraph_runs(paragraph):
    """
    Merge adjacent runs with identical formatting, then fuse into a tagged string.
    Returns (valid_runs, fused_text) or (None, None) if no translatable text.
    """
    raw_runs = list(paragraph.runs)
    if not raw_runs:
        return None, None

    merged_runs = merge_adjacent_runs_inplace(raw_runs, _can_merge_adjacent_pptx_runs)
    valid_runs, fused_text = fuse_merged_paragraph_runs(merged_runs)
    if not valid_runs:
        return None, None

    return valid_runs, fused_text


def _collect_text_units(prs):
    """
    Walk all slides and collect translatable text units.
    Each unit: {"paragraph": p, "valid_runs": [...], "fused_text": "..."}
    """
    units = []

    def process_text_frame(text_frame):
        for paragraph in text_frame.paragraphs:
            valid_runs, fused_text = _fuse_paragraph_runs(paragraph)
            if valid_runs:
                units.append({
                    "paragraph": paragraph,
                    "valid_runs": valid_runs,
                    "fused_text": fused_text,
                })

    def process_table(table):
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    process_text_frame(cell.text_frame)

    def process_shapes(shapes):
        for shape in _collect_shapes(shapes):
            if shape.has_text_frame:
                process_text_frame(shape.text_frame)
            if shape.has_table:
                process_table(shape.table)

    for slide in prs.slides:
        process_shapes(slide.shapes)

        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                process_text_frame(slide.notes_slide.notes_text_frame)
        except Exception:
            pass

    return units


def _reassemble_translations(units, translated_texts, target_lang):
    """Write translated text back into the original run objects."""
    for unit, new_fused_text in zip(units, translated_texts):
        valid_runs = unit["valid_runs"]
        new_fused_text = normalize_run_tags(new_fused_text)

        parsed_runs = {}
        for match in re.finditer(r"<r(\d+)>(.*?)</r\1>", new_fused_text, re.DOTALL):
            parsed_runs[int(match.group(1))] = match.group(2)

        if len(parsed_runs) == len(valid_runs):
            for idx, run, leading, trailing, _ in valid_runs:
                new_core = parsed_runs.get(idx, "")
                new_text = f"{leading}{new_core}{trailing}"
                if run.text != new_text:
                    run.text = new_text
            continue

        clean_text = re.sub(r"</?r\d+>", "", new_fused_text)
        _, first_run, first_leading, _, _ = valid_runs[0]
        _, _, _, last_trailing, _ = valid_runs[-1]
        first_run.text = f"{first_leading}{clean_text.strip()}{last_trailing}"
        for _, run, _, _, _ in valid_runs[1:]:
            run.text = ""


def translate_presentation(
    input_path: str,
    output_path: str,
    target_lang: str,
    backend: str = "google",
    api_key: str = "",
    progress: TranslationProgress = None,
    domain_context: str = "",
    telemetry: dict | None = None,
) -> str:
    """
    Translate a .pptx presentation while preserving all non-text elements.
    """
    if progress is None:
        progress = TranslationProgress()
    progress.reset()
    progress.set_phase("Preparing presentation...")

    job_started_at = time.perf_counter()
    phrases = get_protected_phrases()
    telemetry = telemetry if telemetry is not None else {}
    telemetry.setdefault("timings", {})
    telemetry.setdefault("document", {})

    load_started_at = time.perf_counter()
    prs = Presentation(input_path)
    telemetry["timings"]["load_presentation_seconds"] = round(time.perf_counter() - load_started_at, 3)

    telemetry["document"].update({
        "document_mode": "presentation",
        "backend": backend,
        "slide_count": len(prs.slides),
    })

    progress.set_phase("Collecting text from slides...")
    collect_started_at = time.perf_counter()
    units = _collect_text_units(prs)
    telemetry["timings"]["collect_pptx_seconds"] = round(time.perf_counter() - collect_started_at, 3)

    fused_texts = [u["fused_text"] for u in units]
    non_empty = [t for t in fused_texts if t.strip()]

    estimated_batches = (len(non_empty) + 11) // 12
    total_steps = estimated_batches + FINALIZATION_PROGRESS_STEPS
    progress.set_total(max(1, total_steps))

    telemetry["document"].update({
        "presentation_text_unit_count": len(units),
        "presentation_non_empty_count": len(non_empty),
        "estimated_total_steps": total_steps,
    })

    if not fused_texts:
        progress.set_phase("No translatable text found.")
        progress.advance("Complete")
        prs.save(output_path)
        telemetry["timings"]["total_job_seconds"] = round(time.perf_counter() - job_started_at, 3)
        return output_path

    progress.set_phase("Translating presentation...")
    translation_started_at = time.perf_counter()

    translated_texts = legacy_batch_translate(
        fused_texts,
        target_lang,
        backend,
        api_key,
        phrases,
        progress,
        domain_context,
        telemetry=telemetry,
        skip_decimal_normalization=True,
    )

    telemetry["timings"]["translate_pptx_seconds"] = round(time.perf_counter() - translation_started_at, 3)

    progress.set_phase("Reassembling presentation...")
    _reassemble_translations(units, translated_texts, target_lang)

    progress.set_phase("Saving presentation...")
    save_started_at = time.perf_counter()
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    telemetry["timings"]["save_presentation_seconds"] = round(time.perf_counter() - save_started_at, 3)
    telemetry["timings"]["total_job_seconds"] = round(time.perf_counter() - job_started_at, 3)

    progress.advance("Translation complete")
    return output_path
